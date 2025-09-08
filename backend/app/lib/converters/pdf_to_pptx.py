from pptx import Presentation
from PyPDF2 import PdfReader

def pdf_to_pptx(input_path: str, output_path: str) -> str:
    prs = Presentation()
    reader = PdfReader(input_path)

    for page in reader.pages:
        text = page.extract_text() or ""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        textbox = slide.shapes.placeholders[1]
        textbox.text = text

    prs.save(output_path)
    return output_path
